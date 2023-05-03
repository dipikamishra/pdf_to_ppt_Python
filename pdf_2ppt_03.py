#this code extratct all the text from pdf
import os
from PyPDF2 import PdfFileReader
from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfReader

# Set the path to the PDF file you want to convert
pdf_path = 'EScore-REAlignment20230503_Dipika.pdf'

# Use PyPDF2 to read the PDF file and extract its content
pdf_reader = PdfReader(open(pdf_path, 'rb'))
num_pages = len(pdf_reader.pages)
# Get the first page of the PDF file
first_page = pdf_reader.pages[0]
content = ''
for i in range(num_pages):
    content += pdf_reader.pages[i] #.extract_text()

# Set the path to the PowerPoint file you want to create
pptx_path = 'output.pptx'

# Create a new PowerPoint presentation
prs = Presentation()

# Add a new slide to the presentation
slide = prs.slides.add_slide(prs.slide_layouts[1])

# Add the content of the PDF file to the slide
textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
textbox.text = content

# Save the PowerPoint presentation to disk
prs.save(pptx_path)
